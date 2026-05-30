import os
import sys


def docx_to_pdf(docx_path: str, pdf_path: str = None) -> str:
    from tools.doc_tools import word_to_pdf
    return word_to_pdf(docx_path, pdf_path)


def pdf_to_docx(pdf_path: str, docx_path: str = None) -> str:
    from tools.doc_tools import pdf_to_word
    return pdf_to_word(pdf_path, docx_path)


def xlsx_to_pdf(xlsx_path: str, pdf_path: str = None) -> str:
    """Excel → PDF: pandas → pymupdf Story (headless). Qt fallback. fpdf2 last resort."""
    if pdf_path is None:
        pdf_path = os.path.splitext(xlsx_path)[0] + ".pdf"

    # ── Primary: pandas → HTML → pymupdf Story ──────────────────────────────
    try:
        import pandas as pd
        import fitz

        xlsx     = pd.ExcelFile(xlsx_path)
        all_html = []
        for sheet_name in xlsx.sheet_names:
            df   = pd.read_excel(xlsx, sheet_name=sheet_name).fillna("")
            html = df.to_html(index=False, border=1)
            all_html.append(f"<h2>Sheet: {sheet_name}</h2>{html}")

        combined_html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>
  body  {{font-family:Arial,sans-serif;font-size:9pt;margin:0;}}
  table {{border-collapse:collapse;width:100%;margin-bottom:20px;}}
  th,td {{border:1px solid #666;padding:5px;text-align:left;}}
  th    {{background:#eee;font-weight:bold;}}
  h2    {{color:#2c3e50;border-bottom:1px solid #ccc;padding-bottom:4px;}}
</style></head><body>{"".join(all_html)}</body></html>"""

        story    = fitz.Story(html=combined_html)
        mediabox = fitz.paper_rect("a4-l")        # landscape A4
        margin   = 40
        where    = mediabox + (margin, margin, -margin, -margin)

        writer = fitz.DocumentWriter(pdf_path)
        more   = True
        while more:
            device, more = writer.begin_page(mediabox)
            more, _ = story.place(where)
            story.draw(device)
            writer.end_page()
        writer.close()
        return pdf_path

    except AttributeError:
        pass   # pymupdf < 1.23 — Story not available
    except Exception as e:
        print(f"[xlsx_to_pdf] primary failed: {e}")

    # ── Fallback: pandas → HTML → Qt QPdfWriter ─────────────────────────────
    try:
        import pandas as pd
        from PySide6.QtGui import QTextDocument, QPdfWriter, QPageLayout, QPageSize
        from PySide6.QtCore import QMarginsF
        from PySide6.QtWidgets import QApplication

        if not QApplication.instance():
            _app = QApplication(sys.argv)

        xlsx     = pd.ExcelFile(xlsx_path)
        all_html = []
        for sheet_name in xlsx.sheet_names:
            df   = pd.read_excel(xlsx, sheet_name=sheet_name).fillna("")
            html = df.to_html(index=False, border=1)
            all_html.append(f"<h2>Sheet: {sheet_name}</h2>{html}")

        combined_html = f"""<html><head><style>
          body{{font-family:sans-serif;font-size:9pt;margin:20px;}}
          table{{border-collapse:collapse;width:100%;margin-bottom:25px;}}
          th,td{{border:1px solid #666;padding:6px;text-align:left;}}
          th{{background-color:#eee;font-weight:bold;}}
          h2{{color:#2c3e50;border-bottom:1px solid #ccc;padding-bottom:5px;}}
        </style></head><body>{"".join(all_html)}</body></html>"""

        doc    = QTextDocument()
        doc.setHtml(combined_html)
        writer = QPdfWriter(pdf_path)
        writer.setPageSize(QPageSize(QPageSize.PageSizeId.A4))
        writer.setPageOrientation(QPageLayout.Orientation.Landscape)
        writer.setPageMargins(QMarginsF(10, 10, 10, 10), QPageLayout.Unit.Millimeter)
        doc.print_(writer)
        return pdf_path

    except Exception as e:
        print(f"[xlsx_to_pdf] Qt fallback failed: {e}")

    # ── Last resort: openpyxl → fpdf2 plain text ────────────────────────────
    try:
        import openpyxl
        from fpdf import FPDF

        wb  = openpyxl.load_workbook(xlsx_path)
        ws  = wb.active
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Courier", size=8)
        for row in ws.iter_rows(values_only=True):
            line = " | ".join(str(c) if c is not None else "" for c in row)
            pdf.multi_cell(0, 5, line)
        pdf.output(pdf_path)
        return pdf_path
    except Exception as e:
        raise RuntimeError(f"xlsx_to_pdf: all paths failed — {e}")


def pdf_to_xlsx(pdf_path: str, xlsx_path: str = None) -> str:
    """PDF → Excel: pdfplumber table extraction. Fallback: pymupdf text rows."""
    if xlsx_path is None:
        xlsx_path = os.path.splitext(pdf_path)[0] + ".xlsx"

    # ── Primary: pdfplumber ──────────────────────────────────────────────────
    try:
        import pdfplumber
        import pandas as pd

        sheets_written = 0
        with pdfplumber.open(pdf_path) as pdf:
            with pd.ExcelWriter(xlsx_path, engine="openpyxl") as writer:
                for i, page in enumerate(pdf.pages):
                    table = page.extract_table()
                    if table:
                        df = pd.DataFrame(table[1:], columns=table[0])
                        df.to_excel(writer, sheet_name=f"Page_{i+1}", index=False)
                        sheets_written += 1
        if sheets_written > 0:
            return xlsx_path
        print("[pdf_to_xlsx] pdfplumber found no tables, trying text fallback…")
    except Exception as e:
        print(f"[pdf_to_xlsx] pdfplumber failed: {e}")

    # ── Fallback: pymupdf text rows → one sheet per page ────────────────────
    try:
        import fitz
        import openpyxl

        wb  = openpyxl.Workbook()
        wb.remove(wb.active)  # remove default blank sheet
        pdf = fitz.open(pdf_path)

        for page_num in range(len(pdf)):
            page = pdf[page_num]
            ws   = wb.create_sheet(title=f"Page_{page_num+1}")
            row_idx = 1
            for block in page.get_text("blocks"):
                text = block[4].strip()
                if text:
                    for line in text.splitlines():
                        if line.strip():
                            ws.cell(row=row_idx, column=1, value=line.strip())
                            row_idx += 1

        pdf.close()
        wb.save(xlsx_path)
        return xlsx_path

    except Exception as e:
        print(f"[pdf_to_xlsx] pymupdf fallback failed: {e}")

    # ── Last resort: pypdf ───────────────────────────────────────────────────
    try:
        import openpyxl
        from pypdf import PdfReader

        wb     = openpyxl.Workbook()
        ws     = wb.active
        reader = PdfReader(pdf_path)
        row    = 1
        for page in reader.pages:
            text = page.extract_text() or ""
            for line in text.splitlines():
                if line.strip():
                    ws.cell(row=row, column=1, value=line.strip())
                    row += 1
        wb.save(xlsx_path)
        return xlsx_path
    except Exception as e:
        raise RuntimeError(f"pdf_to_xlsx: all paths failed — {e}")


def pptx_to_pdf(pptx_path: str, pdf_path: str = None) -> str:
    """PPTX → PDF: slide text → HTML → pymupdf Story. Qt & fpdf2 fallbacks."""
    if pdf_path is None:
        pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"

    def _build_pptx_html(pptx_path):
        from pptx import Presentation
        prs      = Presentation(pptx_path)
        all_html = []
        for i, slide in enumerate(prs.slides):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            parts.append(f"<p>{t}</p>")
            all_html.append(f"<div class='slide'><h2>Slide {i+1}</h2>{''.join(parts)}</div>")
        return f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>
  body{{font-family:Arial,sans-serif;margin:0;}}
  .slide{{page-break-after:always;border:1px solid #ccc;padding:20px;margin-bottom:10px;}}
  h2{{color:#2980b9;}} p{{margin:4px 0;}}
</style></head><body>{"".join(all_html)}</body></html>"""

    # ── Primary: pymupdf Story ───────────────────────────────────────────────
    try:
        import fitz
        html     = _build_pptx_html(pptx_path)
        story    = fitz.Story(html=html)
        mediabox = fitz.paper_rect("a4")
        margin   = 40
        where    = mediabox + (margin, margin, -margin, -margin)

        writer = fitz.DocumentWriter(pdf_path)
        more   = True
        while more:
            device, more = writer.begin_page(mediabox)
            more, _ = story.place(where)
            story.draw(device)
            writer.end_page()
        writer.close()
        return pdf_path

    except AttributeError:
        pass
    except Exception as e:
        print(f"[pptx_to_pdf] pymupdf failed: {e}")

    # ── Fallback: Qt QPdfWriter ─────────────────────────────────────────────
    try:
        from PySide6.QtGui import QTextDocument, QPdfWriter, QPageLayout, QPageSize
        from PySide6.QtCore import QMarginsF
        from PySide6.QtWidgets import QApplication

        if not QApplication.instance():
            _app = QApplication(sys.argv)

        html   = _build_pptx_html(pptx_path)
        doc    = QTextDocument()
        doc.setHtml(html)
        writer = QPdfWriter(pdf_path)
        writer.setPageSize(QPageSize(QPageSize.PageSizeId.A4))
        writer.setPageMargins(QMarginsF(15, 15, 15, 15), QPageLayout.Unit.Millimeter)
        doc.print_(writer)
        return pdf_path

    except Exception as e:
        print(f"[pptx_to_pdf] Qt fallback failed: {e}")

    # ── Last resort: fpdf2 ──────────────────────────────────────────────────
    try:
        from pptx import Presentation
        from fpdf import FPDF

        prs = Presentation(pptx_path)
        pdf = FPDF()
        pdf.set_font("Helvetica", size=10)
        for slide in prs.slides:
            pdf.add_page()
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            pdf.multi_cell(0, 6, t)
        pdf.output(pdf_path)
        return pdf_path
    except Exception as e:
        raise RuntimeError(f"pptx_to_pdf: all paths failed — {e}")


def pdf_to_pptx(pdf_path: str, pptx_path: str = None) -> str:
    """PDF → PPTX using pymupdf text extraction (better than pypdf)."""
    if pptx_path is None:
        pptx_path = os.path.splitext(pdf_path)[0] + ".pptx"

    try:
        import fitz
        from pptx import Presentation
        from pptx.util import Inches

        prs            = Presentation()
        prs.slide_width  = Inches(13.333)
        prs.slide_height = Inches(7.5)
        pdf = fitz.open(pdf_path)

        for page in pdf:
            text  = page.get_text("text").strip()
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            txBox = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5),
                prs.slide_width - Inches(1),
                prs.slide_height - Inches(1),
            )
            tf             = txBox.text_frame
            tf.word_wrap   = True
            tf.text        = text

        pdf.close()
        prs.save(pptx_path)
        return pptx_path

    except Exception as e:
        # Fallback: pypdf
        from pptx import Presentation
        from pptx.util import Inches
        from pypdf import PdfReader

        prs            = Presentation()
        prs.slide_width  = Inches(13.333)
        prs.slide_height = Inches(7.5)
        reader = PdfReader(pdf_path)

        for page in reader.pages:
            text  = page.extract_text() or ""
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5),
                prs.slide_width - Inches(1),
                prs.slide_height - Inches(1),
            )
            tf           = txBox.text_frame
            tf.word_wrap = True
            tf.text      = text

        prs.save(pptx_path)
        return pptx_path


def csv_to_xlsx(csv_path: str, xlsx_path: str = None, delimiter: str = ",") -> str:
    import pandas as pd
    if xlsx_path is None:
        xlsx_path = os.path.splitext(csv_path)[0] + ".xlsx"
    df = pd.read_csv(csv_path, sep=delimiter)
    df.to_excel(xlsx_path, index=False)
    return xlsx_path


def xlsx_to_csv(xlsx_path: str, csv_path: str = None) -> str:
    import pandas as pd
    if csv_path is None:
        csv_path = os.path.splitext(xlsx_path)[0] + ".csv"
    df = pd.read_excel(xlsx_path)
    df.to_csv(csv_path, index=False)
    return csv_path


def image_to_pdf(image_paths: list, pdf_path: str = "output.pdf") -> str:
    import img2pdf
    with open(pdf_path, "wb") as f:
        f.write(img2pdf.convert(image_paths))
    return pdf_path


def text_to_pdf(text: str, pdf_path: str = "output.pdf") -> str:
    from tools.doc_tools import text_to_pdf as _t2p
    return _t2p(text, pdf_path)


def merge_pdfs(pdf_list: list, output_path: str = "merged.pdf") -> str:
    from tools.doc_tools import merge_pdfs as _merge
    return _merge(pdf_list, output_path)


def split_pdf(pdf_path: str, output_dir: str = None) -> list:
    from tools.doc_tools import split_pdf as _sp
    return _sp(pdf_path, output_dir)


def pdf_to_image(pdf_path: str, output_dir: str = None, fmt: str = "png", dpi: int = 200) -> list:
    import fitz
    if output_dir is None:
        output_dir = os.path.dirname(pdf_path) or "."
    os.makedirs(output_dir, exist_ok=True)
    doc   = fitz.open(pdf_path)
    base  = os.path.splitext(os.path.basename(pdf_path))[0]
    paths = []
    mat   = fitz.Matrix(dpi / 72, dpi / 72)
    for i, page in enumerate(doc):
        pix      = page.get_pixmap(matrix=mat)
        out_path = os.path.join(output_dir, f"{base}_page_{i+1}.{fmt}")
        pix.save(out_path)
        paths.append(out_path)
    doc.close()
    return paths


def pdf_to_text(pdf_path: str, output_path: str = None) -> str:
    if output_path is None:
        output_path = os.path.splitext(pdf_path)[0] + ".txt"
    try:
        import fitz
        doc  = fitz.open(pdf_path)
        text = "".join(page.get_text() for page in doc)
        doc.close()
    except Exception:
        from pypdf import PdfReader
        reader = PdfReader(pdf_path)
        text   = "".join(p.extract_text() or "" for p in reader.pages)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)
    return output_path
